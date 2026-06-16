#!/usr/bin/env python
"""生成 Swin-Unet 代码讲解 PPT——面向无 ML 基础的同学。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1a, 0x1a, 0x1a)
WHITE = RGBColor(0xff, 0xff, 0xff)
BLUE = RGBColor(0x3a, 0x8e, 0xd8)
YELLOW = RGBColor(0xf0, 0xc0, 0x40)
GREEN = RGBColor(0x3c, 0xb3, 0x71)
RED = RGBColor(0xe7, 0x4c, 0x3c)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xbb, 0xbb, 0xbb)
BG = RGBColor(0x1a, 0x23, 0x33)
CODE_BG = RGBColor(0x0d, 0x11, 0x17)
CODE_FG = RGBColor(0xcc, 0xcc, 0xcc)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def title(slide, text, top=Inches(0.3), size=Pt(36)):
    tb = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = WHITE
    p.font.bold = True


def body(slide, text, top=Inches(1.5), left=Inches(0.8), size=Pt(22),
         color=WHITE, width=Inches(11.5)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def big_text(slide, text, top=Inches(2.0), size=Pt(36), color=WHITE):
    tb = slide.shapes.add_textbox(Inches(1.0), top, Inches(11.3), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def code_box(slide, code, top, height, size=Pt(17)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), top, Inches(12.3), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    for i, line in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = CODE_FG
        p.font.name = 'Consolas'
        p.space_after = Pt(2)


def two_col(slide, left, right, top=Inches(1.6), size=Pt(20)):
    t1 = body(slide, left, top=top, left=Inches(0.5), size=size, width=Inches(5.8))
    t2 = body(slide, right, top=top, left=Inches(6.8), size=size, width=Inches(5.8))
    return t1, t2


def box(slide, text, left, top, w, h, bg_color, text_color=WHITE, size=Pt(18)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '30 行代码看懂 Swin Transformer', top=Inches(2.2), size=Pt(44))
body(s, '—— CT 图像器官分割 · 从零理解 · 不需要 ML 基础', top=Inches(3.5), size=Pt(22), color=GRAY)

# ════════════════════════════════════════════
# Slide 2 — What does this program do?
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '这个程序在做什么？', size=Pt(36))
two_col(s,
        '输入\n\n┌──────────────┐\n│              │\n│  512 × 512   │\n│  黑白 CT 图   │\n│              │\n└──────────────┘\n\n一张腹部 CT 切片\n只有不同深浅的灰色',
        '输出\n\n┌──────────────┐\n│ 🟥🟥  🟩🟩🟩  │\n│ 🟥🟥  🟩🟩🟩  │\n│    🟦🟦🟦    │\n│  🟧🟧🟧🟧🟧   │\n│  🟧🟧🟧🟧🟧   │\n└──────────────┘\n\n每个像素标了颜色\n红=主动脉 绿=胆囊\n蓝=肾 橙=肝 …')
body(s, '任务：给每个像素打上器官标签 — 这叫"医学图像分割"', top=Inches(5.2), size=Pt(18), color=GRAY)

# ════════════════════════════════════════════
# Slide 3 — The core problem
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '核心难题：怎么让程序"看懂"一张图？', size=Pt(34))
body(s,
         '老方法 CNN（卷积神经网络）：用一个 3×3 的小窗口扫过全图\n\n'
         '    窗口太小 → 只看到附近像素 → 不知道"这片灰色是肝还是胃"\n'
         '    → 需要很多层才能看到远处 → 慢、容易漏\n\n'
         'Transformer 的思路：让每个像素都"看到"全图所有像素\n\n'
         '    但 224×224 = 50,176 个像素\n'
         '    每个像素看所有其他人 = 50,176² ≈ 25 亿次计算 → 太慢了！\n\n'
         '🪟 Swin Transformer 的答案：把图切成小块，只在块内"互看"\n'
         '   然后通过"移位"让块之间也能交流 → 又快又好',
         top=Inches(1.5), size=Pt(20))

# ════════════════════════════════════════════
# Slide 4 — Window Partition (the key idea)
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '核心思想 ①：把图切成 7×7 的窗口', size=Pt(34))
body(s,
         '224×224 的图 → 切成 32×32 = 1024 个窗口 → 每个窗口 7×7 = 49 像素\n\n'
         '┌────┬────┬────┬────┐\n'
         '│ W1 │ W2 │ W3 │ ·· │    每个窗口内：49 个像素互相"看"\n'
         '├────┼────┼────┼────┤    窗口之间：暂不交流\n'
         '│ W5 │ W6 │ W7 │ ·· │\n'
         '├────┼────┼────┼────┤    计算量：1024 × 49² ≈ 250 万\n'
         '│ ·· │ ·· │ ·· │ ·· │    vs 全图 50,176² = 25 亿\n'
         '└────┴────┴────┴────┘    → 快了 1000 倍！\n\n'
         '类比：50 人的教室 → 分成 7 人一组，先组内讨论',
         top=Inches(1.5), size=Pt(19))

# ════════════════════════════════════════════
# Slide 5 — Shifted Window (why it's called "Swin")
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '核心思想 ②：移位窗口 — 让窗口之间也能交流', size=Pt(34))
two_col(s,
        '偶数层：正常窗口\n\n┌────┬────┬────┐\n│ A  │ B  │ C  │\n├────┼────┼────┤\n│ D  │ E  │ F  │\n├────┼────┼────┤\n│ G  │ H  │ I  │\n└────┴────┴────┘\n\nA 和 B 分在不同窗口\n→ 它们的像素永远不会交流\n→ 信息被窗口边界锁死 🔒',
        '奇数层：移位 3 像素再切\n\n┌─┬────┬──┬─┐\n│A│  B │ C│D│\n├─┼────┼──┼─┤\n│E│  F │ G│H│\n├─┼────┼──┼─┤\n│I│  J  │K │L│\n└─┴────┴──┴─┘\n\n原来的 A 和 B 的像素\n→ 现在在同一个新窗口里\n→ 可以交流了！🔓')
body(s, '类比：先和同桌讨论 → 换座位 → 再和新同桌讨论。两轮下来，全班信息都流动起来了。',
     top=Inches(5.3), size=Pt(18), color=YELLOW)

# ════════════════════════════════════════════
# Slide 6 — Swin Block = W-MSA + SW-MSA
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '一个 Swin Block = 偶数层 + 奇数层', size=Pt(34))
code_box(s,
         '# 这就是整个架构的"心脏" — 其他都是包装\n'
         'class SwinTransformerBlock(nn.Module):\n'
         '    def forward(self, x):\n'
         '        # 偶数层: W-MSA  — 正常窗口，窗口内互看\n'
         '        # 奇数层: SW-MSA — 移位 3px 后窗口，跨窗口交流\n'
         '        shortcut = x                    # 记住输入（抄近道）\n'
         '        x = self.norm(x)                # 标准化\n'
         '        x = self.attention(x)           # 🎯 核心：注意力计算\n'
         '        x = x + shortcut                # 加上抄的近道\n'
         '        shortcut = x                    # 再记一次\n'
         '        x = self.norm(x)\n'
         '        x = self.mlp(x)                 # 学习更复杂的模式\n'
         '        x = x + shortcut                # 又加上近道\n'
         '        return x',
         top=Inches(1.3), height=Inches(5.0), size=Pt(16))
body(s, '关键参数：窗口=7×7 像素 | 移位=3 像素 | 每个 Block=1个 W-MSA + 1个 SW-MSA',
     top=Inches(6.4), size=Pt(16), color=GRAY)

# ════════════════════════════════════════════
# Slide 7 — Window partition code
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '"窗口划分"在代码里长什么样？', size=Pt(34))
code_box(s,
         'def window_partition(x, window_size=7):\n'
         '    """把整张图切成 7×7 窗口"""\n'
         '    B, H, W, C = x.shape        # 批大小, 高, 宽, 通道\n'
         '    # 把 H×W 的图 reshape 成 (H/7)×(W/7) 个 7×7 小窗口\n'
         '    x = x.reshape(B, H//7, 7, W//7, 7, C)\n'
         '    # 合并维度 → 每行是一个窗口的 49 个像素\n'
         '    windows = x.permute(0, 1, 3, 2, 4, 5).reshape(-1, 7*7, C)\n'
         '    return windows   # (B × num_windows, 49, C)',
         top=Inches(1.5), height=Inches(3.0), size=Pt(18))
code_box(s,
         '# 反过来：把窗口拼回图\n'
         'def window_reverse(windows, window_size=7, H=56, W=56):\n'
         '    """把窗口还原成完整图片"""\n'
         '    B = int(windows.shape[0] / (H*W/7/7))\n'
         '    x = windows.reshape(B, H//7, W//7, 7, 7, -1)\n'
         '    x = x.permute(0, 1, 3, 2, 4, 5).reshape(B, H, W, -1)\n'
         '    return x',
         top=Inches(4.8), height=Inches(2.5), size=Pt(18))

# ════════════════════════════════════════════
# Slide 8 — Inside a window: Attention
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '窗口内部："注意力"在做什么？（用比喻理解）', size=Pt(32))
body(s,
         '窗口里有 7×7 = 49 个像素。每个像素都要问自己一个问题：\n\n'
         '  "我应该关注窗口里哪些其他像素？"\n\n'
         '  类比 — 在教室里看 CT 片：\n'
         '  • 看到一片灰色 → 不确定是肝还是胃\n'
         '  • 往左看看（找肝的特征）→ 往右看看（找胃的特征）\n'
         '  • 综合判断 → "这片灰色的形状更像肝" → 关注左侧\n\n'
         '  代码做的事完全一样：\n'
         '  每个像素计算自己和窗口内所有其他像素的"相似度"\n'
         '  → 相似度高的给大权重（"多关注"）\n'
         '  → 相似度低的给小权重（"少关注"）\n'
         '  → 加权求和，更新自己的"理解"',
         top=Inches(1.5), size=Pt(20))

# ════════════════════════════════════════════
# Slide 9 — Attention code (minimal)
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '注意力计算的 4 步（代码核心）', size=Pt(34))
code_box(s,
         '# WindowAttention.forward() — 窗口内注意力\n'
         'def forward(self, x, mask=None):\n'
         '    # 第 1 步：每个像素生成 3 个向量\n'
         '    # Q=我想找什么  K=我能提供什么  V=我实际有什么\n'
         '    Q = self.qkv(x)[:,:,:self.dim]        # Query\n'
         '    K = self.qkv(x)[:,:,self.dim:self.dim*2]  # Key\n'
         '    V = self.qkv(x)[:,:,self.dim*2:]          # Value\n\n'
         '    # 第 2 步：算相似度 → Q 和 K 的内积\n'
         '    attn = (Q @ K.transpose(-2, -1)) / sqrt(d)  # 缩放\n\n'
         '    # 第 3 步：转成 [0,1] 的权重 → Softmax\n'
         '    attn = softmax(attn, dim=-1)\n\n'
         '    # 第 4 步：加权求和 → 每个像素更新自己的理解\n'
         '    x = attn @ V\n'
         '    return self.proj(x)',
         top=Inches(1.3), height=Inches(5.8), size=Pt(16))

# ════════════════════════════════════════════
# Slide 10 — QKV analogy
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, 'QKV 是什么意思？（用人话版）', size=Pt(34))
body(s,
         '把每个像素想象成一个"人"，现在他们要开会讨论"这是哪个器官"：\n\n'
         '┌──────────────┬────────────────────────────────┐\n'
         '│ Q = Query    │ "我想知道什么？"                  │\n'
         '│              │ 比如："我是肝脏区域吗？"           │\n'
         '├──────────────┼────────────────────────────────┤\n'
         '│ K = Key      │ "我能提供什么信息？"               │\n'
         '│              │ 比如："我长得很像肝的形状！"        │\n'
         '├──────────────┼────────────────────────────────┤\n'
         '│ V = Value    │ "我的实际样子是什么？"              │\n'
         '│              │ 比如："我是偏暗、偏圆的灰色区域"     │\n'
         '├──────────────┼────────────────────────────────┤\n'
         '│ Q @ K^T      │ 问(Q)和答(K)的匹配度 → 权重         │\n'
         '├──────────────┼────────────────────────────────┤\n'
         '│ attn @ V     │ 按权重汇总大家的"样子" → 综合判断    │\n'
         '└──────────────┴────────────────────────────────┘',
         top=Inches(1.5), size=Pt(18))

# ════════════════════════════════════════════
# Slide 11 — U-Net structure
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, 'Swin Block 怎么搭成 U 型网络？', size=Pt(34))
body(s,
         'Swin Block 是最小的积木。把积木按规律排列，就搭出了 U 型网络：\n\n'
         '  编码器（往下压缩，提取特征）    解码器（往上放大，恢复图片）\n'
         '  ┌─────────────────────┐    ┌─────────────────────┐\n'
         '  │ 阶段1  2×Swin Block  │ ─→ │ 阶段1  2×Swin Block  │ ← 恢复细节\n'
         '  │  ↓ 分辨率减半         │    │  ↑ 分辨率加倍         │\n'
         '  │ 阶段2  2×Swin Block  │ ─→ │ 阶段2  2×Swin Block  │ ← 恢复形状\n'
         '  │  ↓ 分辨率减半         │    │  ↑ 分辨率加倍         │\n'
         '  │ 阶段3  6×Swin Block  │ ─→ │ 阶段3  2×Swin Block  │ ← 恢复关系\n'
         '  │  ↓ 分辨率减半         │    │  ↑ 分辨率加倍         │\n'
         '  │ 阶段4  2×Swin Block  │ ─→ │ 阶段4  2×Swin Block  │ │\n'
         '  │         瓶颈层          │    │            ↑          │ │\n'
         '  └─────────────────────┘    └───────────╪──────────┘ │\n'
         '                             跳跃连接 ──┘  (保留细节)',
         top=Inches(1.5), size=Pt(17))

# ════════════════════════════════════════════
# Slide 12 — File map (from easy to hard)
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '代码导航：按理解难度排序', size=Pt(34))
body(s,
         '⭐  trainer.py        训练循环 — 最好理解，从这开始\n'
         '     "喂数据 → 算错多少 → 调整参数 → 重复"\n\n'
         '⭐⭐  train.py / test.py  入口 — 设定参数，启动程序\n\n'
         '⭐⭐⭐  networks/vision_transformer.py  组装 — Swin Block 拼成 U 网\n\n'
         '⭐⭐⭐⭐  networks/swin_transformer_unet_skip_expand_decoder_sys.py\n'
         '        Swin Block 本身 + 窗口划分 + 注意力 — 核心中的核心\n\n'
         '⭐⭐⭐⭐  utils.py  Dice Loss — 判断"对错"的标准\n\n'
         '📖  建议阅读顺序：trainer.py → vision_transformer.py → swin_transformer_unet',
         top=Inches(1.5), size=Pt(20))

# ════════════════════════════════════════════
# Slide 13 — Why this approach
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, 'Swin Transformer 好在哪？', size=Pt(34))
two_col(s,
        'vs 传统 CNN（卷积网络）\n\n❌ 3×3 窗口太小\n     → 只能看到附近\n❌ 要堆很多层\n     → 训练慢、容易漏\n\n✅ 7×7 窗口更大\n✅ 窗口移位 → 全局信息流动\n✅ 每个像素知道全图上下文',
        'vs 普通 Transformer\n\n❌ 全图 50,176² 次计算\n     → 显存放不下\n     → 根本不可行\n\n✅ 窗口内 49² 次计算\n✅ 线性复杂度 → 显存放得下\n✅ 移位机制 → 仍有全局视角')
body(s, '一句话：取了 CNN（局部高效）和 Transformer（全局理解）的折中',
     top=Inches(5.5), size=Pt(18), color=YELLOW)

# ════════════════════════════════════════════
# Slide 14 — Results
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '实际效果怎么样？', size=Pt(34))
body(s,
         '测试 12 个新的 CT 扫描（训练时从未见过）：\n\n'
         '  ✅  肝脏 DSC 0.929  → 大器官，效果最好\n'
         '  ✅  脾脏 DSC 0.876  → 大器官\n'
         '  ✅  肾脏 DSC ~0.81  → 中等器官\n'
         '  🟡 胃   DSC 0.732  → 形状多变\n'
         '  ⚠️  胆囊 DSC 0.611  → 小器官，困难\n'
         '  ⚠️  胰腺 DSC 0.541  → 最小器官，最难\n\n'
         '  平均 DSC 0.761 | 原论文 DSC 0.791（差距来自 GPU 显存限制）\n\n'
         '  DSC 含义：0 = 完全不重合   1 = 完美重合   0.93 = 93% 的区域对了',
         top=Inches(1.5), size=Pt(20))

# ════════════════════════════════════════════
# Slide 15 — Three key numbers
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '记住 3 个数字就理解了 Swin-Unet', size=Pt(36))
body(s,
         '7    窗口大小 = 7×7 像素\n'
         '     为什么是 7？太小 → 信息不足   太大 → 计算变慢   7 是实验最优\n\n'
         '3    移位距离 = 3 像素\n'
         '     为什么是 3？约窗口大小的一半   恰好让新旧窗口重叠最多\n\n'
         '49   每个窗口内的像素数 = 7² = 49\n'
         '     49² 次计算 × 1024 个窗口 ≈ 250 万    vs 全图 25 亿\n'
         '     这就是它"快"的秘密',
         top=Inches(1.5), size=Pt(24))

# ════════════════════════════════════════════
# Slide 16 — Summary
# ════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title(s, '如果只记 3 句话', size=Pt(42))
body(s,
         '1.  把图切成 7×7 小块\n'
         '    每个块内的像素互相"看"，理解局部是什么\n\n'
         '2.  奇数层窗口移位 3 像素\n'
         '    原来隔壁的像素变成同窗 → 信息跨窗口流动 → 全局理解\n\n'
         '3.  交替堆叠 → 形成 Swin Block\n'
         '    先缩小再放大 → U 型网络 → 输出每个像素的器官标签\n\n\n'
         '代码入口：trainer.py（训练）→ vision_transformer.py（组装）\n'
         '          → swin_transformer_unet...py（🏠 核心：Swin Block）',
         top=Inches(1.5), size=Pt(24))

prs.save('Swin-Unet/代码讲解.pptx')
print(f'Saved: 代码讲解.pptx  ({len(prs.slides)} slides)')
